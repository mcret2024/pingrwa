# -*- coding: utf-8 -*-
"""
실리카STO 운영 매뉴얼 V2 — 디자인 빌더
Output: D:\OneDrive\바탕 화면\실리카STO_V2_designed.pptx
"""
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============================================================
# Color Palette — Navy + Gold (premium RWA / blockchain feel)
# ============================================================
NAVY        = RGBColor(0x1E, 0x27, 0x61)
NAVY_DARK   = RGBColor(0x0F, 0x17, 0x35)
GOLD        = RGBColor(0xF5, 0xC4, 0x00)
GOLD_DIM    = RGBColor(0xCD, 0xA3, 0x00)
CORAL       = RGBColor(0xFF, 0x6B, 0x35)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG    = RGBColor(0xF4, 0xF6, 0xFA)
CARD_BG     = RGBColor(0xF8, 0xFA, 0xFC)
BORDER      = RGBColor(0xE5, 0xE7, 0xEB)
TEXT_DARK   = RGBColor(0x1F, 0x29, 0x37)
TEXT_MUTED  = RGBColor(0x6B, 0x72, 0x80)
RED         = RGBColor(0xDC, 0x26, 0x26)
RED_BG      = RGBColor(0xFE, 0xE2, 0xE2)
GREEN       = RGBColor(0x10, 0xB9, 0x81)
BLUE        = RGBColor(0x3B, 0x82, 0xF6)
PURPLE      = RGBColor(0x7C, 0x3A, 0xED)

# ============================================================
# Fonts (Windows-safe Korean fonts)
# ============================================================
FONT_HEADER = '나눔스퀘어'  # Korean display feel; fallback handled by PPT
FONT_BODY   = 'Malgun Gothic'
FONT_MONO   = 'Consolas'

# Slide dimensions: 16:9
SLIDE_W = 13.333  # inches
SLIDE_H = 7.5

# ============================================================
prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_rect(slide, x, y, w, h, fill, line=None, line_w=0.5, corner=0):
    """Add rectangle (or rounded if corner > 0)."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner > 0 else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    if corner > 0:
        try:
            sh.adjustments[0] = corner
        except Exception:
            pass
    return sh

def add_text(slide, x, y, w, h, text, *, size=14, font=FONT_BODY, bold=False,
             color=TEXT_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.3):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb

def add_multi_text(slide, x, y, w, h, paragraphs, *, default_size=14,
                   default_color=TEXT_DARK, line_spacing=1.4, align=PP_ALIGN.LEFT):
    """paragraphs: list of dicts {text, size?, bold?, color?, font?, indent?}"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = item.get('align', align)
        p.line_spacing = item.get('line_spacing', line_spacing)
        if item.get('indent'):
            p.level = item['indent']
        r = p.add_run()
        r.text = item['text']
        r.font.name = item.get('font', FONT_BODY)
        r.font.size = Pt(item.get('size', default_size))
        r.font.bold = item.get('bold', False)
        r.font.color.rgb = item.get('color', default_color)
    return tb

def add_circle(slide, x, y, d, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

# ============================================================
# Common decorators — chip header + page number + footer
# ============================================================
def header(slide, section, slide_no, total):
    # Top accent strip
    add_rect(slide, 0, 0, SLIDE_W, 0.06, NAVY)
    # Section chip
    chip = add_rect(slide, 0.55, 0.32, 1.8, 0.34, NAVY, corner=0.3)
    add_text(slide, 0.55, 0.31, 1.8, 0.36, section.upper(),
             size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    # Page number
    add_text(slide, SLIDE_W - 1.6, 0.32, 1.0, 0.36,
             f'{slide_no:02d} / {total:02d}',
             size=10, color=TEXT_MUTED, font=FONT_MONO, align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

def footer(slide):
    # Footer line
    add_rect(slide, 0.55, SLIDE_H - 0.4, SLIDE_W - 1.1, 0.02, BORDER)
    # Brand
    add_text(slide, 0.55, SLIDE_H - 0.35, SLIDE_W - 1.1, 0.3,
             'SILICA STO  ·  운영 매뉴얼 V2  ·  rwa.silicachainholding.com',
             size=9, color=TEXT_MUTED, anchor=MSO_ANCHOR.MIDDLE)

def title_block(slide, title, subtitle=None, y=0.95):
    add_text(slide, 0.55, y, SLIDE_W - 1.1, 0.7, title,
             size=34, color=NAVY, bold=True, font=FONT_HEADER, line_spacing=1.1)
    if subtitle:
        add_text(slide, 0.55, y + 0.78, SLIDE_W - 1.1, 0.4, subtitle,
                 size=14, color=TEXT_MUTED, line_spacing=1.2)
    # Gold underline accent
    add_rect(slide, 0.55, y + 0.65 + (0.5 if subtitle else 0), 0.7, 0.06, GOLD)

# ============================================================
# Slide 1 — Cover
# ============================================================
def slide_cover():
    s = blank()
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY_DARK)
    # Gold corner accent
    add_rect(s, 0, 0, 0.4, SLIDE_H, GOLD)
    add_rect(s, SLIDE_W - 0.06, 0, 0.06, SLIDE_H, GOLD)
    # Section chip
    add_text(s, 1.2, 1.6, 6, 0.4, '◆ RWA · STO · BLOCKCHAIN',
             size=12, color=GOLD, bold=True, font=FONT_HEADER)
    # Main title
    add_text(s, 1.2, 2.2, 10, 1.3, 'SILICA STO',
             size=76, color=WHITE, bold=True, font=FONT_HEADER, line_spacing=1.0)
    add_text(s, 1.2, 3.5, 10, 0.7, '운영 매뉴얼 V2',
             size=32, color=GOLD, bold=True, font=FONT_HEADER, line_spacing=1.0)
    # Description
    add_text(s, 1.2, 4.5, 10, 0.6,
             'Silica 광산 자산 STO 플랫폼 · 운영자 가이드',
             size=16, color=WHITE, line_spacing=1.3)
    # Tech tags
    add_text(s, 1.2, 5.2, 10, 0.4,
             'Solana · Phantom Wallet · didit.me KYC · Hostinger',
             size=11, color=GOLD_DIM, font=FONT_MONO, line_spacing=1.3)
    # Bottom URL
    add_text(s, 1.2, 6.6, 10, 0.3,
             'rwa.silicachainholding.com',
             size=11, color=GOLD, font=FONT_MONO)

# ============================================================
# Section accent (small chip used inside slides)
# Helper: card with title + body
# ============================================================
def card(slide, x, y, w, h, fill=CARD_BG, border=BORDER, line_w=0.5):
    return add_rect(slide, x, y, w, h, fill, line=border, line_w=line_w, corner=0.05)

def info_card(slide, x, y, w, h, title, body, *, accent=BLUE, title_size=14, body_size=12):
    card(slide, x, y, w, h)
    add_rect(slide, x, y, 0.08, h, accent)  # left accent bar
    add_text(slide, x + 0.25, y + 0.18, w - 0.4, 0.35, title,
             size=title_size, bold=True, color=TEXT_DARK)
    add_text(slide, x + 0.25, y + 0.55, w - 0.4, h - 0.7, body,
             size=body_size, color=TEXT_MUTED, line_spacing=1.5)

def step_circle(slide, x, y, num, color=NAVY):
    add_circle(slide, x, y, 0.5, color)
    add_text(slide, x, y, 0.5, 0.5, str(num),
             size=18, bold=True, color=WHITE, font=FONT_HEADER,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

def caution_box(slide, x, y, w, h, text, *, label='주의', bg=RED_BG, accent=RED):
    add_rect(slide, x, y, w, h, bg, line=accent, line_w=1, corner=0.05)
    add_text(slide, x + 0.25, y + 0.15, w - 0.5, 0.35,
             f'⚠ {label}', size=12, bold=True, color=accent)
    add_text(slide, x + 0.25, y + 0.5, w - 0.5, h - 0.6, text,
             size=12, color=TEXT_DARK, line_spacing=1.5)

def tip_box(slide, x, y, w, h, text, *, label='참고'):
    add_rect(slide, x, y, w, h, RGBColor(0xDB, 0xEA, 0xFE),
             line=BLUE, line_w=1, corner=0.05)
    add_text(slide, x + 0.25, y + 0.15, w - 0.5, 0.35,
             f'ℹ {label}', size=12, bold=True, color=BLUE)
    add_text(slide, x + 0.25, y + 0.5, w - 0.5, h - 0.6, text,
             size=12, color=TEXT_DARK, line_spacing=1.5)

# ============================================================
# Total page count (excluding cover)
# ============================================================
TOTAL = 35

# ============================================================
# Slide 2 (Orig 1) — KYC (1/4)
# ============================================================
def s_kyc_1():
    s = blank()
    header(s, 'KYC · 1/4', 1, TOTAL)
    title_block(s, 'KYC 인증 개요', 'didit.me 외부 시스템을 통한 신원 인증')

    # Left: 4 key points
    points = [
        ('KYC 미진행 시', 'USDT 입금만 가능 (투자/거래/스테이킹 등은 제한)', RED),
        ('투자 자격 요건', 'KYC 사전 진행 필수', NAVY),
        ('외부 시스템 사용', 'https://didit.me/', BLUE),
        ('기본 무료 한도', '월 500건까지 무료, 초과 시 유료', GREEN),
    ]
    y = 2.4
    for i, (title, body, color) in enumerate(points):
        info_card(s, 0.55, y, 6.0, 0.95, title, body, accent=color)
        y += 1.1

    # Right: KYC flow visual
    card(s, 7.0, 2.4, 5.78, 4.3)
    add_text(s, 7.2, 2.55, 5.5, 0.4, 'KYC 진행 흐름',
             size=14, bold=True, color=NAVY)
    flow_steps = [
        ('1', '지갑 연결', 'Phantom 모바일/PC'),
        ('2', '기본 정보 입력', '이름 · 생년월일'),
        ('3', '신분증 + 셀카', 'didit.me 카메라'),
        ('4', '자동 승인', '4개 요인 통과 시'),
    ]
    sy = 3.1
    for num, t, sub in flow_steps:
        step_circle(s, 7.3, sy, num)
        add_text(s, 7.95, sy, 4.7, 0.3, t, size=13, bold=True, color=TEXT_DARK)
        add_text(s, 7.95, sy + 0.28, 4.7, 0.3, sub, size=11, color=TEXT_MUTED)
        sy += 0.85

    footer(s)

# ============================================================
# Slide 3 (Orig 2) — KYC (2/4)
# ============================================================
def s_kyc_2():
    s = blank()
    header(s, 'KYC · 2/4', 2, TOTAL)
    title_block(s, '인증 요인 및 검토', '4개 요인 자동 검증 + 수동 보완')

    # 4 verification factors
    add_text(s, 0.55, 2.4, 6, 0.4, '인증 요인 (4개)',
             size=13, bold=True, color=NAVY)
    factors = [
        ('신분증', 'OCR + 발급기관 검증', '🪪'),
        ('얼굴 셀카', 'Liveness 인증', '📷'),
        ('얼굴 매칭', '신분증 사진 ↔ 셀카 유사도', '👥'),
        ('데이터 검증', '글로벌 DB 대조', '🗄'),
    ]
    fy = 2.9
    for label, desc, icon in factors:
        card(s, 0.55, fy, 6.0, 0.75)
        add_text(s, 0.85, fy + 0.15, 0.6, 0.5, icon, size=22,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        add_text(s, 1.5, fy + 0.1, 4.4, 0.3, label, size=14, bold=True, color=TEXT_DARK)
        add_text(s, 1.5, fy + 0.42, 4.4, 0.3, desc, size=11, color=TEXT_MUTED)
        fy += 0.9

    # Right: Manual review
    info_card(s, 7.0, 2.4, 5.78, 1.5,
              '시스템 반려 시',
              '관리자가 didit.me 대시보드에서 확인 후\n수동으로 승인 또는 반려 처리',
              accent=GOLD)

    caution_box(s, 7.0, 4.05, 5.78, 1.4,
                '대부분의 반려 사유: 신분증 사진과 얼굴 셀카의 유사도 부족',
                label='흔한 반려 케이스')

    tip_box(s, 7.0, 5.55, 5.78, 1.2,
            '4개 요인 모두 통과 시 자동 KYC 승인 처리',
            label='자동 승인')

    footer(s)

# ============================================================
# Slide 4 (Orig 3) — KYC (3/4)
# ============================================================
def s_kyc_3():
    s = blank()
    header(s, 'KYC · 3/4', 3, TOTAL)
    title_block(s, '거부 시 재신청', '관리자가 거부한 경우 유저 재신청 흐름')

    # Center flow
    card(s, 0.55, 2.4, SLIDE_W - 1.1, 4.4)
    add_text(s, 0.85, 2.6, SLIDE_W - 1.7, 0.4,
             '거부 → 재신청 흐름',
             size=15, bold=True, color=NAVY)

    flow = [
        ('관리자 검토', 'didit.me 대시보드에서 결과 확인', '👤'),
        ('거부 결정', '유사도 부족 등 사유로 거부', '✕'),
        ('유저 통지', '거부 사유 + 신청일/거부일 표시', '📢'),
        ('재신청', '유저가 Resubmit 버튼 클릭', '🔄'),
        ('재인증 진행', '새 didit 세션으로 다시 시작', '✓'),
    ]
    fx = 0.85
    fy = 3.4
    box_w = 2.34
    gap = 0.15
    for i, (title, desc, icon) in enumerate(flow):
        x = fx + i * (box_w + gap)
        card(s, x, fy, box_w, 1.4, fill=WHITE, border=BORDER)
        add_text(s, x, fy + 0.15, box_w, 0.5, icon, size=24,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        add_text(s, x, fy + 0.7, box_w, 0.3, title, size=12, bold=True,
                 color=TEXT_DARK, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.05, fy + 1.0, box_w - 0.1, 0.4, desc, size=10,
                 color=TEXT_MUTED, align=PP_ALIGN.CENTER, line_spacing=1.3)
        # Arrow between
        if i < len(flow) - 1:
            arrow_x = x + box_w + 0.01
            add_text(s, arrow_x, fy + 0.55, 0.13, 0.3, '▶', size=12,
                     color=GOLD, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

    tip_box(s, 0.85, 5.1, SLIDE_W - 1.7, 1.4,
            '유저가 거부됨을 확인하면 KYC 페이지에서 Resubmit 버튼이 자동 노출됩니다.\n'
            '이름/생년월일이 불일치한 경우 → "Edit Info" 버튼 (기본 정보 페이지로 이동)\n'
            'KYC 제공자 크레딧 부족 → "Contact Support" (관리자 문의)',
            label='재신청 액션 분기')

    footer(s)

# ============================================================
# Slide 5 (Orig 4) — KYC (4/4)
# ============================================================
def s_kyc_4():
    s = blank()
    header(s, 'KYC · 4/4', 4, TOTAL)
    title_block(s, '유사도 조정', '잦은 수동 검토를 줄이는 didit 설정')

    card(s, 0.55, 2.4, SLIDE_W - 1.1, 4.4)
    add_text(s, 0.85, 2.6, SLIDE_W - 1.7, 0.5,
             '신분증 얼굴 ↔ 셀카 유사도 임계값 조정',
             size=18, bold=True, color=NAVY)
    add_text(s, 0.85, 3.15, SLIDE_W - 1.7, 0.4,
             '잦은 검토가 불편한 경우 didit.me 대시보드에서 임계값을 낮춰주세요',
             size=12, color=TEXT_MUTED)

    # 3 columns: too high / too low / recommended
    cols = [
        ('임계값 높음', '엄격한 검증', '실제 본인도 자주 반려됨\n수동 검토 빈도 증가', RED),
        ('임계값 낮음', '관대한 검증', '검증 정확도 하락\n오인 승인 가능성', CORAL),
        ('임계값 적정', '균형 설정', '대부분 자동 통과\n예외만 수동 검토', GREEN),
    ]
    cx = 0.85
    cy = 3.85
    cw = 4.0
    ch = 2.5
    for i, (title, sub, desc, color) in enumerate(cols):
        x = cx + i * (cw + 0.18)
        card(s, x, cy, cw, ch, fill=WHITE, border=BORDER)
        add_rect(s, x, cy, cw, 0.08, color)
        add_text(s, x + 0.2, cy + 0.2, cw - 0.4, 0.4, title,
                 size=15, bold=True, color=color)
        add_text(s, x + 0.2, cy + 0.65, cw - 0.4, 0.3, sub,
                 size=11, color=TEXT_MUTED)
        add_text(s, x + 0.2, cy + 1.1, cw - 0.4, ch - 1.3, desc,
                 size=12, color=TEXT_DARK, line_spacing=1.5)

    footer(s)

# ============================================================
# Slide 6 (Orig 5) — Phantom 셀카 이슈
# ============================================================
def s_kyc_phantom_issue():
    s = blank()
    header(s, 'KYC · 알려진 이슈', 5, TOTAL)
    title_block(s, 'Phantom 익스플로러 셀카 이슈',
                'Phantom 인앱 브라우저에서 발생하는 알려진 현상')

    caution_box(s, 0.55, 2.4, SLIDE_W - 1.1, 1.5,
                'Phantom 익스플로러에서 셀카 촬영 시 동일 메시지가 3회 반복 노출될 수 있습니다.\n'
                'Phantom 측 인앱 브라우저 이슈로 사이트에서 제어 불가합니다.',
                label='알려진 현상')

    card(s, 0.55, 4.1, 6.1, 2.7)
    add_text(s, 0.85, 4.3, 5.7, 0.4, '대응 권장',
             size=14, bold=True, color=NAVY)
    items = [
        '유저에게 같은 메시지가 반복돼도 정상 진행하도록 안내',
        '필요 시 일반 Chrome / Safari / Samsung Internet 으로 전환 권장',
        'KYC 핸드오프 시스템으로 외부 브라우저 전환 옵션 제공 (v793+)',
    ]
    iy = 4.85
    for txt in items:
        add_circle(s, 0.9, iy + 0.05, 0.18, GOLD)
        add_text(s, 1.2, iy, 5.4, 0.3, txt, size=12, color=TEXT_DARK, line_spacing=1.4)
        iy += 0.55

    card(s, 6.85, 4.1, 5.93, 2.7)
    add_text(s, 7.15, 4.3, 5.5, 0.4, '관련 시스템',
             size=14, bold=True, color=NAVY)
    add_multi_text(s, 7.15, 4.8, 5.5, 1.9, [
        {'text': '· KYC Handoff 시스템', 'bold': True, 'size': 12},
        {'text': '  Phantom 인앱 → Chrome 등 외부 브라우저로\n  KYC 세션 안전하게 전환', 'size': 11, 'color': TEXT_MUTED},
        {'text': '· didit.me 사용', 'bold': True, 'size': 12},
        {'text': '  외부 KYC 제공자 측 한계',  'size': 11, 'color': TEXT_MUTED},
    ], default_size=11)

    footer(s)

# ============================================================
# Two-column flow slide template (user flow vs admin flow)
# ============================================================
def two_column_flow(s, title, subtitle, left_title, left_steps, right_title, right_steps,
                    left_color=BLUE, right_color=PURPLE):
    title_block(s, title, subtitle)
    col_w = (SLIDE_W - 1.1 - 0.3) / 2
    cy = 2.4
    ch = SLIDE_H - 2.4 - 0.6

    for idx, (cx, ctitle, steps, ccolor) in enumerate([
        (0.55, left_title, left_steps, left_color),
        (0.55 + col_w + 0.3, right_title, right_steps, right_color),
    ]):
        card(s, cx, cy, col_w, ch)
        add_rect(s, cx, cy, col_w, 0.5, ccolor)
        add_text(s, cx + 0.25, cy, col_w - 0.5, 0.5, ctitle,
                 size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        sy = cy + 0.75
        # Adaptive step spacing so 5+ steps don't overflow the card bottom (6.9").
        n_tuples = sum(1 for st in steps if isinstance(st, tuple))
        n_plain  = len(steps) - n_tuples
        tuple_gap_default = 0.95
        plain_gap_default = 0.7
        available = ch - 0.85  # card height minus title band + bottom padding
        needed = n_tuples * tuple_gap_default + n_plain * plain_gap_default
        scale = min(1.0, available / needed) if needed > 0 else 1.0
        tuple_gap = tuple_gap_default * scale
        plain_gap = plain_gap_default * scale
        for i, step in enumerate(steps, 1):
            step_circle(s, cx + 0.3, sy, i, ccolor)
            # step is (title, desc) or just title
            if isinstance(step, tuple):
                t, d = step
                add_text(s, cx + 0.95, sy + 0.02, col_w - 1.2, 0.3, t,
                         size=12, bold=True, color=TEXT_DARK)
                if d:
                    add_text(s, cx + 0.95, sy + 0.32, col_w - 1.2, 0.5, d,
                             size=10.5, color=TEXT_MUTED, line_spacing=1.4)
                sy += tuple_gap
            else:
                add_text(s, cx + 0.95, sy + 0.1, col_w - 1.2, 0.4, step,
                         size=12, color=TEXT_DARK, line_spacing=1.4)
                sy += plain_gap

# ============================================================
# Slide 7 (Orig 6) — 입금하기 (1/2)
# ============================================================
def s_deposit_1():
    s = blank()
    header(s, '입금 · 1/2', 6, TOTAL)
    two_column_flow(s, '입금하기', '유저 흐름',
        '유저 플로우',
        [
            ('토큰 선택', '지갑 → 입금 → USDT 또는 토큰'),
            ('수량 입력', '입금할 수량 입력'),
            ('보내기 클릭', '지갑 연결 + 전송'),
            ('진행 사항 확인', '입금 진행 상태 실시간 표시'),
            ('관리자 승인', '승인 시 진행사항 제거'),
        ],
        '관련 페이지',
        [
            '/user/deposit.html',
            '/admin/deposits.html',
            '/admin/token-deposits.html',
        ],
        left_color=BLUE, right_color=NAVY)
    footer(s)

# ============================================================
# Slide 8 (Orig 7) — 입금하기 (2/2)
# ============================================================
def s_deposit_2():
    s = blank()
    header(s, '입금 · 2/2', 7, TOTAL)
    two_column_flow(s, '입금 승인', '관리자 흐름',
        '관리자 플로우',
        [
            ('승인대기 목록 확인', 'USDT/토큰 입금에서 대기중 항목'),
            ('승인 처리', '승인 클릭으로 즉시 반영'),
        ],
        '체크 포인트',
        [
            '온체인 입금 정상 수신 여부',
            '입금자 주소 = 등록된 유저 지갑',
            '입금 금액 정확성',
            '메모/태그 (해당 시) 확인',
        ],
        left_color=PURPLE, right_color=GOLD_DIM)
    footer(s)

# ============================================================
# Slide 9 (Orig 8) — 입금거절하기 (1/2)
# ============================================================
def s_deposit_reject_1():
    s = blank()
    header(s, '입금 거절 · 1/2', 8, TOTAL)
    two_column_flow(s, '입금 거절', '관리자 처리 가이드',
        '관리자 플로우',
        [
            ('거절 내용 입력', '왜 거절하는지 사유'),
            ('연락처 입력', '유저가 연락할 수 있는 정보'),
            ('거절 확정', '버튼 클릭으로 처리'),
            ('환불 처리', '정상 입금 확인 시 외부 경로 환불'),
        ],
        '거절 사유 예시',
        [
            '입금자 지갑이 등록 지갑과 불일치',
            '입금 메모/태그 오류',
            '비정상 토큰 또는 마이그레이션 토큰',
            '제재 대상 주소',
        ],
        left_color=RED, right_color=NAVY)
    footer(s)

# ============================================================
# Slide 10 (Orig 9) — 입금거절하기 (2/2)
# ============================================================
def s_deposit_reject_2():
    s = blank()
    header(s, '입금 거절 · 2/2', 9, TOTAL)
    title_block(s, '거절 시 유저 대응', '환불 안내 흐름')

    card(s, 0.55, 2.4, SLIDE_W - 1.1, 4.4)
    add_text(s, 0.85, 2.6, SLIDE_W - 1.7, 0.5,
             '유저 측에서 보이는 안내', size=15, bold=True, color=NAVY)

    add_multi_text(s, 0.85, 3.3, SLIDE_W - 1.7, 3.4, [
        {'text': '거절 사실 + 사유 + 관리자 연락처가 유저 히스토리에 표시',
         'size': 13, 'color': TEXT_DARK, 'bold': True},
        {'text': '', 'size': 6},
        {'text': '유저 행동 흐름',
         'size': 13, 'color': NAVY, 'bold': True},
        {'text': '  1.  거절 알림 확인 (히스토리 페이지)', 'size': 12, 'color': TEXT_DARK},
        {'text': '  2.  관리자가 남긴 연락처로 연락 (이메일/카카오톡 등)', 'size': 12, 'color': TEXT_DARK},
        {'text': '  3.  소명 자료 제출 (입금 증빙, 본인 확인 등)', 'size': 12, 'color': TEXT_DARK},
        {'text': '  4.  관리자 검토 후 정상 입금 확인 시', 'size': 12, 'color': TEXT_DARK},
        {'text': '  5.  외부 경로 (블록체인 환불 송금) 로 처리', 'size': 12, 'color': TEXT_DARK},
        {'text': '', 'size': 6},
        {'text': '※ 환불은 사이트 내 자동 처리가 아니므로 관리자가 직접 송금합니다',
         'size': 11, 'color': RED, 'bold': True},
    ])

    footer(s)

# ============================================================
# Slide 11 (Orig 10) — 출금
# ============================================================
def s_withdraw():
    s = blank()
    header(s, '출금', 10, TOTAL)
    title_block(s, '출금 처리', '본인 지갑으로만 출금 가능')

    # Two cards: User + Admin
    cw = (SLIDE_W - 1.1 - 0.3) / 2
    cy = 2.4
    ch = 4.4

    # Left: User
    card(s, 0.55, cy, cw, ch)
    add_rect(s, 0.55, cy, cw, 0.5, BLUE)
    add_text(s, 0.8, cy, cw - 0.5, 0.5, '유저 흐름', size=14, bold=True,
             color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_multi_text(s, 0.85, cy + 0.75, cw - 0.6, ch - 1, [
        {'text': '출금 요청', 'size': 14, 'bold': True, 'color': TEXT_DARK},
        {'text': '본인 지갑 주소로만 가능 (다른 지갑 X)', 'size': 11, 'color': TEXT_MUTED},
        {'text': '', 'size': 6},
        {'text': '결과 확인', 'size': 14, 'bold': True, 'color': TEXT_DARK},
        {'text': '히스토리 → 출금 항목에서 상태/사유 확인', 'size': 11, 'color': TEXT_MUTED},
        {'text': '', 'size': 6},
        {'text': '반려 시', 'size': 14, 'bold': True, 'color': RED},
        {'text': '관리자가 입력한 반려 사유가 표시됨', 'size': 11, 'color': TEXT_MUTED},
    ])

    # Right: Admin
    cx2 = 0.55 + cw + 0.3
    card(s, cx2, cy, cw, ch)
    add_rect(s, cx2, cy, cw, 0.5, PURPLE)
    add_text(s, cx2 + 0.25, cy, cw - 0.5, 0.5, '관리자 처리',
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_multi_text(s, cx2 + 0.3, cy + 0.75, cw - 0.6, ch - 1, [
        {'text': '선택 전송', 'size': 14, 'bold': True, 'color': GREEN},
        {'text': '관리자가 검토 후 "선택 전송" 버튼으로 처리', 'size': 11, 'color': TEXT_MUTED},
        {'text': '', 'size': 6},
        {'text': '반려', 'size': 14, 'bold': True, 'color': RED},
        {'text': '반려 사유 입력 필수 → 유저 히스토리에 노출', 'size': 11, 'color': TEXT_MUTED},
        {'text': '', 'size': 6},
        {'text': '주의 사항', 'size': 14, 'bold': True, 'color': NAVY},
        {'text': '• 온체인 송금 후 되돌릴 수 없음', 'size': 11, 'color': TEXT_MUTED},
        {'text': '• 주소 정확성 사전 확인 필수', 'size': 11, 'color': TEXT_MUTED},
        {'text': '• 가스 비용 + 수수료 정산 확인', 'size': 11, 'color': TEXT_MUTED},
    ])

    footer(s)

# ============================================================
# Slide 12-13 (Orig 11-12) — 투자
# ============================================================
def s_invest_1():
    s = blank()
    header(s, '투자 · 1/2', 11, TOTAL)
    two_column_flow(s, '투자하기', '유저 흐름',
        '유저 플로우',
        [
            ('USDT 투자', '광산 선택 + 투자 금액 입력'),
            ('계약서 작성', '전자계약서 검토 및 서명'),
            ('진행 안내 팝업', '관리자 서명 대기 중 표시'),
            ('승인 완료 팝업', '관리자 서명 완료 후 자동 알림'),
            ('히스토리 확인', '세부 내용은 히스토리 페이지'),
        ],
        '체크 포인트',
        [
            '관리자 서명 이전엔 "대기중" 상태',
            'STO 토큰은 서명 완료 후 자동 지급',
            '반려 시 USDT 자동 환불',
            'KYC 미인증 시 투자 불가',
        ],
        left_color=BLUE, right_color=GOLD_DIM)
    footer(s)

def s_invest_2():
    s = blank()
    header(s, '투자 · 2/2', 12, TOTAL)
    two_column_flow(s, '투자 승인', '관리자 흐름',
        '관리자 플로우',
        [
            ('회계 → 전자계약서', '관리자 서명 대기 목록'),
            ('투자건 확인', '금액, 광산, 유저 정보 검토'),
            ('서명', '전자 서명 완료'),
            ('승인 클릭', '서명 완료 버튼 → STO 자동 지급'),
        ],
        '반려 처리',
        [
            ('반려 사유 작성', '명확한 거절 사유 입력'),
            ('계약 반려 클릭', '버튼 클릭으로 반려 확정'),
            ('USDT 자동 환불', '유저에게 즉시 반환'),
            ('히스토리 기록', '유저/관리자 양측 기록'),
        ],
        left_color=PURPLE, right_color=RED)
    footer(s)

# ============================================================
# Slide 14-15 (Orig 13-14) — 투자 제한 설정
# ============================================================
def s_invest_limit_1():
    s = blank()
    header(s, '투자 제한 · 1/2', 13, TOTAL)
    title_block(s, '투자 총액 제한', '관리자가 광산별 최대 모금액 설정')

    card(s, 0.55, 2.4, SLIDE_W - 1.1, 4.4)
    add_text(s, 0.85, 2.6, SLIDE_W - 1.7, 0.5,
             '설정 위치 및 동작', size=15, bold=True, color=NAVY)

    add_multi_text(s, 0.85, 3.3, SLIDE_W - 1.7, 3.4, [
        {'text': '관리자 → /admin/assets.html → 광산 선택', 'size': 13, 'color': TEXT_DARK, 'bold': True},
        {'text': '', 'size': 6},
        {'text': '관련 필드', 'size': 13, 'color': NAVY, 'bold': True},
        {'text': '  · target_usdt        목표 모금액 (USDT)', 'size': 12, 'color': TEXT_DARK, 'font': FONT_MONO},
        {'text': '  · raised_usdt        현재 모금 누적 (자동 갱신)', 'size': 12, 'color': TEXT_DARK, 'font': FONT_MONO},
        {'text': '  · min_usdt           최소 투자 금액', 'size': 12, 'color': TEXT_DARK, 'font': FONT_MONO},
        {'text': '  · pending_reserved   대기 중 (계약 서명 전 락업)', 'size': 12, 'color': TEXT_DARK, 'font': FONT_MONO},
        {'text': '', 'size': 6},
        {'text': '제한 동작', 'size': 13, 'color': NAVY, 'bold': True},
        {'text': '  유저들의 총투자 + 대기중 = target_usdt 초과 시 신규 투자 차단', 'size': 12, 'color': TEXT_DARK},
    ])
    footer(s)

def s_invest_limit_2():
    s = blank()
    header(s, '투자 제한 · 2/2', 14, TOTAL)
    title_block(s, '제한 초과 시', '유저에게 표시되는 흐름')

    card(s, 0.55, 2.4, SLIDE_W - 1.1, 4.4)
    add_text(s, 0.85, 2.6, SLIDE_W - 1.7, 0.5,
             '시나리오', size=15, bold=True, color=NAVY)

    scenarios = [
        ('정상 투자 가능', '남은 한도 ≥ 입력 금액', GREEN, '버튼 활성화, 정상 진행'),
        ('한도 임박', '남은 한도 < 입력 금액', GOLD, '경고 + 가능한 최대 금액 안내'),
        ('한도 도달', '남은 한도 = 0 또는 음수', RED, '투자 자체 차단, 모금 종료 안내'),
    ]
    sy = 3.4
    for label, condition, color, action in scenarios:
        card(s, 0.85, sy, SLIDE_W - 1.7, 1.0, fill=WHITE, border=BORDER)
        add_rect(s, 0.85, sy, 0.1, 1.0, color)
        add_text(s, 1.1, sy + 0.13, 3, 0.3, label, size=14, bold=True, color=color)
        add_text(s, 1.1, sy + 0.45, 3, 0.3, condition, size=11, color=TEXT_MUTED, font=FONT_MONO)
        add_text(s, 4.3, sy + 0.3, SLIDE_W - 5.5, 0.5, '→  ' + action,
                 size=12, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
        sy += 1.1
    footer(s)

# ============================================================
# Slide 16-18 (Orig 15-17) — 스테이킹
# ============================================================
def s_staking_1():
    s = blank()
    header(s, '스테이킹 · 1/3', 15, TOTAL)
    title_block(s, '스테이킹 이자 지급', '매월 15일 기준 스테이킹 상태로 판정')

    # Calendar visualization
    card(s, 0.55, 2.4, SLIDE_W - 1.1, 4.4)
    add_text(s, 0.85, 2.6, SLIDE_W - 1.7, 0.5,
             '지급 판정 일정', size=15, bold=True, color=NAVY)

    # Day boxes 13, 14, 15, 16, 17
    days = [
        ('13일', '평상시', GREEN, '스테이킹 / 언스테이킹 가능'),
        ('14일', '잠금 시작', RED, '스테이킹 / 언스테이킹 불가'),
        ('15일', '지급일', GOLD, '이자 자동 지급 (스테이킹 상태 기준)'),
        ('16일', '잠금 유지', RED, '스테이킹 / 언스테이킹 불가'),
        ('17일', '평상시', GREEN, '스테이킹 / 언스테이킹 가능'),
    ]
    dx = 0.85
    dy = 3.4
    dw = (SLIDE_W - 1.7 - 0.4) / 5
    for i, (day, label, color, action) in enumerate(days):
        x = dx + i * (dw + 0.1)
        card(s, x, dy, dw, 2.0, fill=WHITE, border=BORDER)
        add_rect(s, x, dy, dw, 0.6, color)
        add_text(s, x, dy, dw, 0.6, day, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        add_text(s, x, dy + 0.75, dw, 0.4, label, size=13, bold=True, color=color,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + 0.1, dy + 1.2, dw - 0.2, 0.7, action,
                 size=10.5, color=TEXT_DARK, align=PP_ALIGN.CENTER, line_spacing=1.4)
    footer(s)

def s_staking_2():
    s = blank()
    header(s, '스테이킹 · 2/3', 16, TOTAL)
    title_block(s, '잠금 기간 (14~16일)', '3일간 스테이킹/언스테이킹 모두 불가')

    caution_box(s, 0.55, 2.4, SLIDE_W - 1.1, 1.6,
                '매월 14일 ~ 16일 (3일) 동안 스테이킹 / 언스테이킹 진입이 불가합니다.\n'
                '이는 매월 15일 이자 지급 시점의 잔액을 안정적으로 산정하기 위함입니다.',
                label='잠금 기간')

    card(s, 0.55, 4.2, SLIDE_W - 1.1, 2.5)
    add_text(s, 0.85, 4.4, SLIDE_W - 1.7, 0.4,
             '이자 지급 판정 로직', size=14, bold=True, color=NAVY)
    add_multi_text(s, 0.85, 4.9, SLIDE_W - 1.7, 1.7, [
        {'text': '✓  14~16일 내내 스테이킹 상태 유지', 'size': 13, 'color': GREEN, 'bold': True},
        {'text': '   → 정상적으로 이자 지급', 'size': 12, 'color': TEXT_MUTED},
        {'text': '', 'size': 5},
        {'text': '✗  14~16일 중 언스테이킹 된 상태', 'size': 13, 'color': RED, 'bold': True},
        {'text': '   → 해당 월 이자 지급 없음', 'size': 12, 'color': TEXT_MUTED},
    ])
    footer(s)

def s_staking_3():
    s = blank()
    header(s, '스테이킹 · 3/3', 17, TOTAL)
    title_block(s, '이자율 관리', '관리자가 APR 변경, 다음 회차부터 적용')

    two_column_flow(s, '', '',
        '관리자 플로우',
        [
            ('Silica 운영', 'admin → Silica 운영'),
            ('스테이킹 설정', '이자율 (APR) 변경 화면'),
            ('새 APR 입력', '소수점 가능, % 단위'),
            ('저장 클릭', 'apr_history 에 기록'),
        ],
        '적용 정책',
        [
            ('다음 회차부터 적용', '현재 진행 회차는 옛 APR 유지'),
            ('자동 공지', '유저들에게 팝업으로 알림'),
            ('히스토리 보존', 'apr_history 테이블에 이력 누적'),
        ],
        left_color=PURPLE, right_color=GOLD_DIM)
    footer(s)

# ============================================================
# Slide 19-21 (Orig 18-20) — 이자 지급 테스팅
# ============================================================
def s_interest_test_1():
    s = blank()
    header(s, '이자 테스팅 · 1/3', 18, TOTAL)
    title_block(s, '이자 지급 테스팅', '실제 운영 전 테스트용 강제 지급 기능')

    card(s, 0.55, 2.4, SLIDE_W - 1.1, 4.4)
    add_text(s, 0.85, 2.6, SLIDE_W - 1.7, 0.5,
             '관리자 → 실리카운영 → 스테이킹 설정',
             size=15, bold=True, color=NAVY)

    add_multi_text(s, 0.85, 3.3, SLIDE_W - 1.7, 3.4, [
        {'text': '용도', 'size': 13, 'color': NAVY, 'bold': True},
        {'text': '  · 운영 환경에서 실제 이자 지급 흐름을 검증', 'size': 12, 'color': TEXT_DARK},
        {'text': '  · 신규 광산 등록 후 첫 회차 전 사전 확인', 'size': 12, 'color': TEXT_DARK},
        {'text': '  · 특정 사용자에 대한 보상/보정 수동 지급', 'size': 12, 'color': TEXT_DARK},
        {'text': '', 'size': 6},
        {'text': '주의', 'size': 13, 'color': RED, 'bold': True},
        {'text': '  · 같은 월 중복 지급 자동 차단 (per holder duplicate check)', 'size': 12, 'color': TEXT_DARK},
        {'text': '  · 실제 자금이 지급되므로 신중히 사용', 'size': 12, 'color': TEXT_DARK},
        {'text': '  · 테스트 도구는 실제 운영 시 제거 예정', 'size': 12, 'color': TEXT_DARK},
    ])
    footer(s)

def s_interest_test_2():
    s = blank()
    header(s, '이자 테스팅 · 2/3', 19, TOTAL)
    title_block(s, '이자 수령', '스테이킹 / 클레임 페이지에서 수령')

    cols = [
        ('스테이킹 페이지', '/user/staking.html',
         '내 스테이킹 잔액과 누적 이자 한 곳에서 확인 + 클레임'),
        ('클레임 페이지', '/user/claim.html',
         '클레임 가능한 모든 자산 통합 표시 + 일괄 클레임'),
    ]
    cw = (SLIDE_W - 1.1 - 0.3) / 2
    cy = 2.4
    ch = 4.4
    for i, (title, url, desc) in enumerate(cols):
        x = 0.55 + i * (cw + 0.3)
        card(s, x, cy, cw, ch)
        add_rect(s, x, cy, cw, 0.5, NAVY)
        add_text(s, x + 0.25, cy, cw - 0.5, 0.5, title,
                 size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + 0.3, cy + 0.75, cw - 0.6, 0.4, url,
                 size=12, color=GOLD_DIM, font=FONT_MONO)
        add_text(s, x + 0.3, cy + 1.3, cw - 0.6, ch - 1.5, desc,
                 size=13, color=TEXT_DARK, line_spacing=1.6)
    footer(s)

def s_interest_test_3():
    s = blank()
    header(s, '이자 테스팅 · 3/3', 20, TOTAL)
    title_block(s, '수령한 이자 확인', '3개 페이지에서 일관되게 표시')

    pages = [
        ('Staking', '/user/staking.html', '내 스테이킹 현황 + 이자 누적'),
        ('Claim', '/user/claim.html', '미수령 이자 일괄 표시'),
        ('History', '/user/history.html', '수령 이력 시간순 정렬'),
    ]
    cw = (SLIDE_W - 1.1 - 0.4) / 3
    cy = 2.4
    ch = 3.6
    for i, (title, url, desc) in enumerate(pages):
        x = 0.55 + i * (cw + 0.2)
        card(s, x, cy, cw, ch)
        add_rect(s, x, cy, cw, 0.4, GOLD_DIM)
        add_text(s, x, cy, cw, 0.4, title, size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + 0.2, cy + 0.7, cw - 0.4, 0.4, url,
                 size=10, color=BLUE, font=FONT_MONO)
        add_text(s, x + 0.2, cy + 1.3, cw - 0.4, ch - 1.5, desc,
                 size=12, color=TEXT_DARK, line_spacing=1.6)

    tip_box(s, 0.55, 6.20, SLIDE_W - 1.1, 0.85,
            '이자는 소수점 1자리까지만 지급됩니다',
            label='이자 정밀도')
    footer(s)

# ============================================================
# Slide 22-24 (Orig 21-23) — 추천인 보상
# ============================================================
def s_referral_1():
    s = blank()
    header(s, '추천인 · 1/3', 21, TOTAL)
    title_block(s, '추천인 등록', '관리자 → 마케팅 → 추천인 메뉴')

    two_column_flow(s, '', '',
        '관리자 플로우',
        [
            ('마케팅 → 추천인', 'admin → 마케팅 → 추천인 메뉴'),
            ('유저 검색', '추천인이 될 사용자 검색'),
            ('승인 처리', '체크 + 저장으로 권한 부여'),
            ('자동 알림', '대상 유저에게 권한 활성화 안내'),
        ],
        '권한 부여 조건',
        [
            '관리자가 직접 승인한 유저만 추천인 자격',
            '추천 받는 유저 (referee) 측은 제한 없음',
            '동일 유저 동시 추천 불가',
            '추천 수수료율은 광산별 admin 설정',
        ],
        left_color=PURPLE, right_color=GOLD_DIM)
    footer(s)

def s_referral_2():
    s = blank()
    header(s, '추천인 · 2/3', 22, TOTAL)
    title_block(s, '추천인 권한', '활성화 시 추가 메뉴 / 탭 노출')

    card(s, 0.55, 2.4, SLIDE_W - 1.1, 4.4)
    add_text(s, 0.85, 2.6, SLIDE_W - 1.7, 0.5,
             '추천인으로 등록된 유저만 보이는 UI', size=15, bold=True, color=NAVY)

    items = [
        ('레퍼럴 메뉴', 'header & mobile drawer 에 추가', '/user/referral.html'),
        ('레퍼럴 탭 (클레임)', 'Claim 페이지에서 별도 탭 노출', 'referral_bonus kind'),
        ('보상 표시', '추천한 유저들의 투자 누적에 따라 보상 산정', '관리자 설정 % 기반'),
    ]
    iy = 3.3
    for label, desc, sub in items:
        card(s, 0.85, iy, SLIDE_W - 1.7, 1.0, fill=WHITE, border=BORDER)
        add_rect(s, 0.85, iy, 0.1, 1.0, GOLD_DIM)
        add_text(s, 1.1, iy + 0.18, 4, 0.3, label, size=14, bold=True, color=TEXT_DARK)
        add_text(s, 1.1, iy + 0.52, 4, 0.3, desc, size=11, color=TEXT_MUTED)
        add_text(s, 5.3, iy + 0.35, SLIDE_W - 6.5, 0.5, sub,
                 size=11, color=BLUE, font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        iy += 1.1
    footer(s)

def s_referral_3():
    s = blank()
    header(s, '추천인 · 3/3', 23, TOTAL)
    title_block(s, '추천 보상 수령', '클레임 페이지에서 일반 이자와 분리')

    card(s, 0.55, 2.4, SLIDE_W - 1.1, 4.4)
    add_text(s, 0.85, 2.6, SLIDE_W - 1.7, 0.5,
             '수령 흐름', size=15, bold=True, color=NAVY)

    add_multi_text(s, 0.85, 3.3, SLIDE_W - 1.7, 3.4, [
        {'text': '1.  추천한 유저의 투자 완료 시점에 추천 수당 자동 누적', 'size': 13, 'color': TEXT_DARK},
        {'text': '', 'size': 5},
        {'text': '2.  Claim 페이지 → "Referral" 탭에서 누적 수당 확인', 'size': 13, 'color': TEXT_DARK},
        {'text': '', 'size': 5},
        {'text': '3.  "Claim Referral Bonus" 버튼으로 수령', 'size': 13, 'color': TEXT_DARK},
        {'text': '', 'size': 5},
        {'text': '4.  USDT 직접 지급 (지갑 잔액에 반영)', 'size': 13, 'color': TEXT_DARK},
        {'text': '', 'size': 5},
        {'text': '5.  History 페이지에 referral_bonus 항목으로 기록', 'size': 13, 'color': TEXT_DARK},
    ])
    footer(s)

# ============================================================
# Slide 25-26 (Orig 24-25) — 스왑
# ============================================================
def s_swap_1():
    s = blank()
    header(s, '스왑 · 1/2', 24, TOTAL)
    title_block(s, '스왑 (Silica ↔ STO)', 'Silica 토큰을 STO 토큰으로 교환')

    card(s, 0.55, 2.4, 6.2, 4.4)
    add_text(s, 0.85, 2.6, 5.9, 0.5,
             '입력 / 결과', size=15, bold=True, color=NAVY)
    add_multi_text(s, 0.85, 3.3, 5.9, 3.4, [
        {'text': '· STO 토큰: 정수로만 입력', 'size': 13, 'color': TEXT_DARK, 'bold': True},
        {'text': '  소수점 입력 불가', 'size': 11, 'color': TEXT_MUTED},
        {'text': '', 'size': 5},
        {'text': '· 실리카 토큰 수량 자동 표기', 'size': 13, 'color': TEXT_DARK, 'bold': True},
        {'text': '  시세에 따라 지불해야 할 Silica 수량 계산', 'size': 11, 'color': TEXT_MUTED},
        {'text': '', 'size': 5},
        {'text': '· USDT 수수료', 'size': 13, 'color': TEXT_DARK, 'bold': True},
        {'text': '  스왑량 × % (최소 1 USDT)', 'size': 11, 'color': TEXT_MUTED},
    ])

    card(s, 6.95, 2.4, 5.83, 4.4)
    add_text(s, 7.25, 2.6, 5.5, 0.5,
             '스왑 불가 조건', size=15, bold=True, color=RED)
    add_multi_text(s, 7.25, 3.3, 5.5, 3.4, [
        {'text': '✗  실리카 토큰 부족', 'size': 13, 'color': RED, 'bold': True},
        {'text': '   1 STO 이하로 환산되는 양', 'size': 11, 'color': TEXT_MUTED},
        {'text': '', 'size': 5},
        {'text': '✗  USDT 잔액 부족', 'size': 13, 'color': RED, 'bold': True},
        {'text': '   수수료 미지불 시 스왑 불가', 'size': 11, 'color': TEXT_MUTED},
        {'text': '', 'size': 5},
        {'text': '✗  실리카 시세 0 또는 미설정', 'size': 13, 'color': RED, 'bold': True},
        {'text': '   admin → 실리카 시세 설정 필요', 'size': 11, 'color': TEXT_MUTED},
    ])
    footer(s)

def s_swap_2():
    s = blank()
    header(s, '스왑 · 2/2', 25, TOTAL)
    title_block(s, '스왑 수수료 설정', '관리자가 % 기반으로 설정')

    card(s, 0.55, 2.4, SLIDE_W - 1.1, 4.4)
    add_text(s, 0.85, 2.6, SLIDE_W - 1.7, 0.5,
             '설정 → 스왑 수수료', size=15, bold=True, color=NAVY)

    add_multi_text(s, 0.85, 3.3, SLIDE_W - 1.7, 3.4, [
        {'text': '기본 동작', 'size': 13, 'color': NAVY, 'bold': True},
        {'text': '  · 스왑 금액 × 설정 % = 수수료 (USDT)', 'size': 12, 'color': TEXT_DARK},
        {'text': '  · 계산 결과가 1 USDT 미만이어도 1 USDT 부과 (최소 수수료)', 'size': 12, 'color': TEXT_DARK},
        {'text': '', 'size': 6},
        {'text': '특수 케이스: 수수료율 = 0', 'size': 13, 'color': GREEN, 'bold': True},
        {'text': '  · 0% 입력 시 최소 수수료도 발생하지 않음 (완전 무료)', 'size': 12, 'color': TEXT_DARK},
        {'text': '  · 거래 활성화 또는 프로모션 시 활용 가능', 'size': 12, 'color': TEXT_DARK},
        {'text': '', 'size': 6},
        {'text': '의미', 'size': 13, 'color': NAVY, 'bold': True},
        {'text': '  소량 스왑도 최소 1 USDT 부과로 운영 비용 보전', 'size': 12, 'color': TEXT_DARK},
    ])
    footer(s)

# ============================================================
# Slide 27-29 (Orig 26-28) — 거래
# ============================================================
def s_trade_1():
    s = blank()
    header(s, '거래 · 1/3', 26, TOTAL)
    title_block(s, '거래 (P2P 주문판)', '일반 현물 거래소와 동일한 매수/매도 방식')

    card(s, 0.55, 2.4, SLIDE_W - 1.1, 4.4)
    add_text(s, 0.85, 2.6, SLIDE_W - 1.7, 0.5,
             '주문 처리 방식', size=15, bold=True, color=NAVY)

    add_multi_text(s, 0.85, 3.3, SLIDE_W - 1.7, 1.8, [
        {'text': '✓  매수 주문: USDT escrow + 토큰 수령', 'size': 13, 'color': GREEN, 'bold': True},
        {'text': '   주문 시 USDT + 수수료가 escrow 락업', 'size': 11, 'color': TEXT_MUTED},
        {'text': '✓  매도 주문: 토큰 escrow + USDT 수령', 'size': 13, 'color': GREEN, 'bold': True},
        {'text': '   주문 시 토큰만 escrow, 수수료는 체결 시 차감', 'size': 11, 'color': TEXT_MUTED},
    ])

    caution_box(s, 0.85, 5.3, SLIDE_W - 1.7, 1.4,
                '운영 초기 단계로 거래량이 많지 않아 시세 급변동을 완화하는\n'
                '안전장치 (가격 변동 폭 제한 · 서킷브레이커 등) 는 적용되지 않은 상태입니다.',
                label='현재 한계')

    footer(s)

def s_trade_2():
    s = blank()
    header(s, '거래 · 2/3', 27, TOTAL)
    title_block(s, '거래 수수료 정책', '관리자 설정 + 최소 보장 규칙')

    rules = [
        ('수수료율', '관리자 설정 (광산별)', '/admin/assets.html → fee_buyer / fee_seller'),
        ('최소 수수료', '계산값이 0.001 USDT 미만이면 0.001 로 보정', 'applyMinTradeFee 함수'),
        ('USDT 정밀도', '소수점 2자리까지 입력 가능', 'price * 100 정수 검증'),
        ('토큰 정밀도', '정수만 입력 가능 (소수점 X)', 'floor(amount) 처리'),
        ('최소 주문', '1 USDT 미만 주문 불가', 'orderTotal >= 1.0'),
    ]
    sy = 2.4
    for label, value, code in rules:
        card(s, 0.55, sy, SLIDE_W - 1.1, 0.8, fill=WHITE, border=BORDER)
        add_rect(s, 0.55, sy, 0.1, 0.8, GOLD_DIM)
        add_text(s, 0.85, sy + 0.1, 3, 0.3, label, size=13, bold=True, color=TEXT_DARK)
        add_text(s, 0.85, sy + 0.42, 3, 0.3, value, size=11, color=TEXT_MUTED)
        add_text(s, 4.5, sy + 0.27, SLIDE_W - 5.6, 0.4, code,
                 size=11, color=BLUE, font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        sy += 0.9
    footer(s)

def s_trade_3():
    s = blank()
    header(s, '거래 · 3/3', 28, TOTAL)
    title_block(s, '관리자 수수료 관리', '설정 패널에서 광산별 수수료 변경')

    card(s, 0.55, 2.4, SLIDE_W - 1.1, 4.4)
    add_text(s, 0.85, 2.6, SLIDE_W - 1.7, 0.5,
             '변경 흐름', size=15, bold=True, color=NAVY)

    add_multi_text(s, 0.85, 3.3, SLIDE_W - 1.7, 3.4, [
        {'text': '1.  admin → 자산 → 광산 선택', 'size': 13, 'color': TEXT_DARK},
        {'text': '2.  fee_buyer / fee_seller 항목 수정 (% 단위)', 'size': 13, 'color': TEXT_DARK},
        {'text': '3.  저장 → 즉시 신규 주문에 반영', 'size': 13, 'color': TEXT_DARK},
        {'text': '', 'size': 6},
        {'text': '미체결 주문 처리 (v823)', 'size': 13, 'color': NAVY, 'bold': True},
        {'text': '  · 메이커 (등록자) 측: 주문 등록 시점의 fee 유지', 'size': 12, 'color': TEXT_DARK},
        {'text': '  · 테이커 (매칭 시 진입) 측: 현재 fee 적용', 'size': 12, 'color': TEXT_DARK},
        {'text': '  · AMM sweep: 메이커의 placement-time fee 사용 (v823 일관성 fix)', 'size': 12, 'color': TEXT_DARK},
    ])
    footer(s)

# ============================================================
# Slide 30 (Orig 29) — 배당금
# ============================================================
def s_dividend():
    s = blank()
    header(s, '배당금', 29, TOTAL)
    title_block(s, '배당금 실행', '매월 15일 분배, Silica 토큰으로 환산 지급')

    two_column_flow(s, '', '',
        '관리자 플로우',
        [
            ('운영 → 배당금 실행', 'admin → Silica 운영 → 배당 실행'),
            ('분배 가치 입력', 'USDT 기준 분배 총액'),
            ('분기(회차) 선택', '몇 번째 분배인지 선택'),
            ('실행 클릭', '15일 시세로 Silica 환산 → 분배'),
        ],
        '유저 수령',
        [
            ('Claim 페이지', '/user/claim.html 에서 확인'),
            ('Claim 버튼 클릭', '실리카 토큰 수령'),
            ('히스토리 기록', 'dividend kind 로 기록'),
            ('미수령 시 누적', '클레임 안 해도 누적 보존'),
        ],
        left_color=PURPLE, right_color=GOLD_DIM)
    footer(s)

# ============================================================
# Slide 31-32 (Orig 30-31) — 비상 재해
# ============================================================
def s_emergency_1():
    s = blank()
    header(s, '비상 재해 · 1/2', 30, TOTAL)
    title_block(s, '이자/배당금 비상 재해 모드', '14~16일 서버 이슈 발생 시 수동 대응')

    caution_box(s, 0.55, 2.4, SLIDE_W - 1.1, 1.5,
                '아래 긴급 모드는 구성만 되어 있고 실제 테스팅된 기능이 아닙니다.\n'
                '재해/버그로 정상 지급이 안 된 경우 개발자와 함께 데이터 백업 후 진행해주세요.',
                label='주의')

    card(s, 0.55, 4.1, SLIDE_W - 1.1, 2.7)
    add_text(s, 0.85, 4.3, SLIDE_W - 1.7, 0.5,
             '대응 절차', size=15, bold=True, color=NAVY)
    add_multi_text(s, 0.85, 4.85, SLIDE_W - 1.7, 1.9, [
        {'text': '1.  14~16일 사이 서버 이슈 발생 (재해, 오류)', 'size': 12, 'color': TEXT_DARK},
        {'text': '2.  17일에 유저들에게 자동 안내 팝업 표시', 'size': 12, 'color': TEXT_DARK},
        {'text': '3.  유저는 팝업 안내에 따라 언스테이킹 하지 않도록 유지', 'size': 12, 'color': TEXT_DARK},
        {'text': '4.  관리자가 수동 지급 버튼 클릭 → 일괄 지급', 'size': 12, 'color': TEXT_DARK},
        {'text': '5.  지급 완료 후 정상 운영 복귀', 'size': 12, 'color': TEXT_DARK},
    ])

    footer(s)

def s_emergency_2():
    s = blank()
    header(s, '비상 재해 · 2/2', 31, TOTAL)
    title_block(s, '비상 모드 운영 가이드', '데이터 무결성 우선')

    card(s, 0.55, 2.4, SLIDE_W - 1.1, 4.4)
    add_text(s, 0.85, 2.6, SLIDE_W - 1.7, 0.5,
             '운영 원칙', size=15, bold=True, color=NAVY)

    add_multi_text(s, 0.85, 3.3, SLIDE_W - 1.7, 3.4, [
        {'text': '✓  사전 데이터 백업 필수', 'size': 13, 'color': GREEN, 'bold': True},
        {'text': '   holdings / wallet_transactions / interest_claims 전체 export', 'size': 11, 'color': TEXT_MUTED},
        {'text': '', 'size': 5},
        {'text': '✓  중복 지급 방지', 'size': 13, 'color': GREEN, 'bold': True},
        {'text': '   per-holder duplicate check 가 자동 작동', 'size': 11, 'color': TEXT_MUTED},
        {'text': '', 'size': 5},
        {'text': '✓  Disaster Lock', 'size': 13, 'color': GREEN, 'bold': True},
        {'text': '   silicaHasDisasterPending() 가 stake/unstake 차단 (v761)', 'size': 11, 'color': TEXT_MUTED},
        {'text': '', 'size': 5},
        {'text': '✗  운영 정상 시 사용 금지', 'size': 13, 'color': RED, 'bold': True},
        {'text': '   실제 자금 지급이므로 정상 운영 시 호출 X', 'size': 11, 'color': TEXT_MUTED},
    ])
    footer(s)

# ============================================================
# Slide 33-34 (Orig 32-33) — IR 페이지
# ============================================================
def s_ir_1():
    s = blank()
    header(s, 'IR 페이지 · 1/2', 32, TOTAL)
    title_block(s, 'IR 페이지 문서', '관리자 회계 → 자산문서 페이지에서 등록')

    card(s, 0.55, 2.4, SLIDE_W - 1.1, 4.4)
    add_text(s, 0.85, 2.6, SLIDE_W - 1.7, 0.5,
             '문서 카테고리', size=15, bold=True, color=NAVY)

    cats = [
        ('등기부등본', '광산 소유권 / 등기 정보', '🏛'),
        ('자산평가서', '광산 평가 가치 문서', '📊'),
        ('회계자료', '재무제표, 회계 보고서', '📑'),
        ('공식문서', '광산 운영 관련 공식 자료', '🏢'),
        ('증빙문서', '거래/매각 증빙', '🧾'),
        ('일반문서', '기타 첨부 문서', '📎'),
    ]
    cx = 0.85
    cy = 3.3
    cw = (SLIDE_W - 1.7 - 0.4) / 3
    for i, (label, desc, icon) in enumerate(cats):
        col = i % 3
        row = i // 3
        x = cx + col * (cw + 0.2)
        y = cy + row * 1.6
        card(s, x, y, cw, 1.4, fill=WHITE, border=BORDER)
        add_text(s, x + 0.2, y + 0.15, 0.7, 0.7, icon, size=28,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        add_text(s, x + 0.95, y + 0.2, cw - 1.0, 0.35, label,
                 size=13, bold=True, color=TEXT_DARK)
        add_text(s, x + 0.95, y + 0.6, cw - 1.0, 0.7, desc,
                 size=10.5, color=TEXT_MUTED, line_spacing=1.4)
    footer(s)

def s_ir_2():
    s = blank()
    header(s, 'IR 페이지 · 2/2', 33, TOTAL)
    title_block(s, 'IR 페이지 이미지/메타', '/admin/assets.html 에서 등록')

    card(s, 0.55, 2.4, SLIDE_W - 1.1, 4.4)
    add_text(s, 0.85, 2.6, SLIDE_W - 1.7, 0.5,
             '관리자가 등록할 수 있는 항목', size=15, bold=True, color=NAVY)

    items = [
        ('자산 이미지', '광산 hero 이미지 (큰 사진)', 'assets.image_url'),
        ('자산명', '광산 이름 (다국어 지원)', 'assets.name + asset_key_info'),
        ('토큰 로고', 'SilicaSTO 토큰 로고', 'assets.token_image_url'),
        ('지도', 'Google Maps 임베드 URL', 'assets.google_map_url'),
        ('소개 멘트', '광산 소개 / 개요', 'assets.overview + asset_key_info'),
        ('세부 이미지', '여러 장의 갤러리 이미지', 'asset_key_info detail_image_N'),
    ]
    iy = 3.3
    for label, desc, col in items:
        card(s, 0.85, iy, SLIDE_W - 1.7, 0.5, fill=WHITE, border=BORDER)
        add_rect(s, 0.85, iy, 0.08, 0.5, GOLD_DIM)
        add_text(s, 1.05, iy + 0.12, 2.5, 0.3, label, size=12, bold=True, color=TEXT_DARK)
        add_text(s, 3.7, iy + 0.12, 5, 0.3, desc, size=11, color=TEXT_MUTED)
        add_text(s, 8.9, iy + 0.12, SLIDE_W - 9.5, 0.3, col,
                 size=10, color=BLUE, font=FONT_MONO)
        iy += 0.55
    footer(s)

# ============================================================
# Slide 35 (Orig 34) — 운영 중단
# ============================================================
def s_winddown():
    s = blank()
    header(s, '운영 중단', 34, TOTAL)
    title_block(s, '운영 중단 기능', 'Recon RWA 부동산 매각 → Silica 광산 채굴 투자')

    # Top: explanation card
    card(s, 0.55, 2.4, SLIDE_W - 1.1, 1.5)
    add_text(s, 0.85, 2.55, SLIDE_W - 1.7, 0.4,
             'Silica 와 Recon RWA 차이', size=14, bold=True, color=NAVY)
    add_multi_text(s, 0.85, 2.95, SLIDE_W - 1.7, 0.9, [
        {'text': '· 부동산 (Recon): 투자금으로 부동산 매입 → 매각 시 자금 회수 → 매각 기능 사용', 'size': 11, 'color': TEXT_DARK},
        {'text': '· 광산 (Silica): 이미 소유권 보유 광산의 채굴 투자 → 매각 기능 제외 (메뉴 hidden)', 'size': 11, 'color': TEXT_DARK},
    ])

    # Bottom: shutdown features
    card(s, 0.55, 4.1, SLIDE_W - 1.1, 2.7)
    add_text(s, 0.85, 4.25, SLIDE_W - 1.7, 0.4,
             '사이트 정지 기능', size=14, bold=True, color=NAVY)
    features = [
        ('스테이킹 중단', '모든 유저 강제 언스테이킹 + 신규 투자/스테이킹 차단 (출금만 가능)', RED),
        ('서비스 종료', '사이트 최종 종료 → 미출금 유저 정보 엑셀 다운로드 가능', NAVY),
    ]
    fy = 4.75
    for label, desc, color in features:
        card(s, 0.85, fy, SLIDE_W - 1.7, 0.9, fill=WHITE, border=BORDER)
        add_rect(s, 0.85, fy, 0.1, 0.9, color)
        add_text(s, 1.15, fy + 0.15, 3, 0.3, label, size=14, bold=True, color=color)
        add_text(s, 1.15, fy + 0.5, SLIDE_W - 2.2, 0.3, desc,
                 size=11, color=TEXT_MUTED)
        fy += 1.0
    footer(s)

# ============================================================
# Slide 36 (Orig 35) — 외부 서비스
# ============================================================
def s_external():
    s = blank()
    header(s, '외부 서비스', 35, TOTAL)
    title_block(s, '사용 중인 외부 서비스', 'Hostinger · Helius · didit · CoinMarketCap · Yahoo')

    services = [
        ('호스팅', 'Hostinger', '클라우드 호스팅 (향후 트래픽 증가 시 AWS 마이그레이션 검토 가능)', '서버'),
        ('블록체인 API', 'Helius', '월 최소 플랜 49 USD · 입출금/온체인 작업 처리', 'Solana RPC'),
        ('KYC', 'didit.me', '신원 인증 / 신분증 + 셀카 검증 / 월 500건 무료', 'KYC 제공자'),
        ('시세 데이터', '코인마켓캡 무료 API', '상장 코인 시세 조회 (참고용)', '시세'),
        ('환율 API', '야후 파이낸스 무료 API', '원화 환율 조회 (현재 미사용, 확장 대비 보존)', '환율'),
    ]
    sy = 2.4
    for label, brand, desc, tag in services:
        card(s, 0.55, sy, SLIDE_W - 1.1, 0.85, fill=WHITE, border=BORDER)
        add_rect(s, 0.55, sy, 0.1, 0.85, GOLD_DIM)
        add_text(s, 0.85, sy + 0.1, 2.5, 0.3, label, size=12, bold=True, color=NAVY)
        add_text(s, 0.85, sy + 0.42, 2.5, 0.3, brand, size=13, bold=True, color=TEXT_DARK)
        add_text(s, 3.6, sy + 0.18, SLIDE_W - 5.3, 0.6, desc,
                 size=11, color=TEXT_MUTED, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4)
        add_text(s, SLIDE_W - 1.6, sy + 0.3, 0.95, 0.3, tag,
                 size=10, color=GOLD_DIM, font=FONT_MONO, align=PP_ALIGN.RIGHT)
        sy += 0.92
    footer(s)

# ============================================================
# Build all
# ============================================================
slide_cover()
s_kyc_1()
s_kyc_2()
s_kyc_3()
s_kyc_4()
s_kyc_phantom_issue()
s_deposit_1()
s_deposit_2()
s_deposit_reject_1()
s_deposit_reject_2()
s_withdraw()
s_invest_1()
s_invest_2()
s_invest_limit_1()
s_invest_limit_2()
s_staking_1()
s_staking_2()
s_staking_3()
s_interest_test_1()
s_interest_test_2()
s_interest_test_3()
s_referral_1()
s_referral_2()
s_referral_3()
s_swap_1()
s_swap_2()
s_trade_1()
s_trade_2()
s_trade_3()
s_dividend()
s_emergency_1()
s_emergency_2()
s_ir_1()
s_ir_2()
s_winddown()
s_external()

output = r'D:\OneDrive\바탕 화면\실리카STO_V2_designed.pptx'
prs.save(output)
print(f'Saved: {output}')
print(f'Total slides: {len(prs.slides)}')
